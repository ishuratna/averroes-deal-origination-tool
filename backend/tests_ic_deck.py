#!/usr/bin/env python3
"""
The IC screening deck: template fill + the code-only financials grid.

The one property that matters most: the template is the real Blink deck, so
EVERY content shape must be overwritten - any 'Blink' text surviving in the
output means a mapping gap and another company's IC pack leaking a prior
deal's details.
"""
import os
import sys
import warnings

warnings.filterwarnings("ignore")
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
os.environ.setdefault("GCP_PROJECT_ID", "averroes-deal-origination")

from services.ic_memo_deck import financials_grid, render_deck  # noqa: E402

fails = 0


def chk(label, got, want=True):
    global fails
    ok = got == want
    print(("PASS" if ok else "FAIL"), label, "" if ok else f"-> {got!r} (wanted {want!r})")
    if not ok:
        fails += 1


print("── financials_grid: numbers only from the record ──")
c = {"revenue_y1": 5_200_000, "revenue_y1_date": "2025-03-31",
     "revenue_y2": 4_000_000, "revenue_y2_date": "2024-03-31",
     "gross_profit_y1": 4_160_000, "profit_y1": 520_000}
g = financials_grid(c)
chk("13 rows like the template", len(g), 13)
chk("year labels from filing dates", g[1][1:3], ["FY24", "FY25"])
chk("revenue formatted", g[2][1:3], ["£4.00m", "£5.20m"])
chk("growth computed, not guessed", g[3][2], "+30.0%")
chk("gross margin computed", g[6][2], "80.0%")
chk("PBT from profit_y1", g[11][2], "£520k")
chk("EBITDA is n/d, never invented", g[8][1:3], ["n/d", "n/d"])
chk("columns beyond our data are '-'", g[2][4], "-")
chk("empty record never crashes", len(financials_grid({})), 13)
chk("no-data revenue row is all '-'", financials_grid({})[2][1:], ["-"] * 9)

print()
print("── render_deck: every Blink shape overwritten ──")
content = {
    "screening": {"tag": "Proprietary outreach - materials testing SaaS", "rating": "Strong",
                  "note": "Origination-stage screen."},
    "business_description": [{"lead": "Founded", "text": "in 2018, Cambridge."}],
    "alignment": {k: {"verdict": f"{k} verdict", "tone": t} for k, t in
                  [("business_quality", "good"), ("revenue_quality", "warn"),
                   ("market_dynamics", "bad"), ("management", "good"),
                   ("deal_structuring", "warn"), ("sector_focus", "good")]},
    "market_overview": [{"lead": "ICP", "text": "UK labs (AI research)"}],
    "thesis": ["Sticky vertical software."],
    "value_creation": [{"lead": "Pricing", "text": "move to platform fees."}],
    "deal_structuring": ["Founder liquidity discussed."],
    "risks": [{"category": "MANAGEMENT", "risk": "Key-person CEO", "mitigant": "Verify bench in DD"}],
    "exit": {"buyer_types": ["Strategic - instrument OEMs"], "examples": "Instron (AI research)"},
    "questions": [{"lead": "Churn", "text": "What is gross churn by cohort?"}],
}
company = {"name": "Plastometrex", "sector": "Materials testing", "revenue_y1": 5_200_000}
data = render_deck(company, content, g)
chk("returns real pptx bytes", data[:2] == b"PK")

import io
from pptx import Presentation
prs = Presentation(io.BytesIO(data))
all_text = []
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame:
            all_text.append(sh.text_frame.text)
        if sh.has_table:
            for r in sh.table.rows:
                for cell in r.cells:
                    all_text.append(cell.text)
blob = "\n".join(all_text)
chk("no Blink text survives anywhere", "Blink" not in blob and "Weiser" not in blob and "Fidelity" not in blob)
chk("company name on every slide", blob.count("Plastometrex") >= 4)
chk("rating chip says Strong", "Strong" in blob)
chk("financials landed in the table", "£5.20m" in blob)
chk("risk + mitigant written", "Key-person CEO" in blob and "Verify bench in DD" in blob)
chk("research tag preserved", "(AI research)" in blob)
chk("question lead written, no manual numbering (template auto-numbers)",
    "Churn" in blob and "1. Churn" not in blob)
chk("fixed alignment labels kept", "Business Quality" in blob and "Sector Focus" in blob)
chk("unused risk rows blanked (no leftover categories)", blob.count("MARKET DYNAMICS") <= 1)

print()
print(f"{fails} FAILURES" if fails else "ALL PASS")
sys.exit(1 if fails else 0)
