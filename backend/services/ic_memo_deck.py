"""
IC memo DECK: the 4-slide screening pack in the Blink CIM format (per Ishu,
27 Aug 2026). The Blink deck itself is the template - templates/
ic_memo_deck_template.pptx - and every content shape is OVERWRITTEN by shape
id, so the format is followed strictly and no Blink text can survive.

Slides: 1 Deal Screening Summary, 2 Overview (description / alignment chips /
key financials / market), 3 Investment Rationale (thesis / value creation /
deal structuring / risks table / exit), 4 Follow-up Questions for Diligence.

Rules, in the order they matter:
  * NUMBERS come only from the stored record. The financials table is filled
    in CODE from CH-filed figures; the model never writes a number into it.
    Missing cells read "n/d" (not disclosed), never a guess.
  * The model may research market/competitor context with Google Search, and
    every claim that came from research (not the record) is tagged
    "(AI research)" so the IC can see what is sourced vs company-stated.
  * One grounded Gemini call per deck, through the shared budget (weight 2 -
    it is doing real research). No em dashes anywhere.
"""
import copy
import io
import json
import logging
import os
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                             "..", "templates", "ic_memo_deck_template.pptx")

# The six alignment criteria are FIXED by the template; the model writes only
# the one-line verdict and a tone for each.
ALIGNMENT_KEYS = [("business_quality", "Business Quality"),
                  ("revenue_quality", "Revenue Quality"),
                  ("market_dynamics", "Market Dynamics"),
                  ("management", "Management"),
                  ("deal_structuring", "Deal Structuring"),
                  ("sector_focus", "Sector Focus")]

_TONE_FILL = {"good": (0xE2, 0xEF, 0xDA), "warn": (0xFF, 0xF2, 0xCC), "bad": (0xF8, 0xCB, 0xAD)}
_RATING_FILL = {"Strong": (0x2E, 0x7D, 0x32), "Mixed": (0xD9, 0x77, 0x06), "Weak": (0xB9, 0x1C, 0x1C)}


# ── The financials table: CODE ONLY, from the stored record ──────────────────

def _gbp(v) -> str:
    if v in (None, ""):
        return "n/d"
    try:
        v = float(v)
    except (TypeError, ValueError):
        return "n/d"
    if abs(v) >= 1_000_000:
        return f"£{v / 1_000_000:.2f}m"
    return f"£{v / 1_000:.0f}k"


def _yr(datestr) -> str:
    s = str(datestr or "")
    for tok in s.replace("-", " ").split():
        if tok.isdigit() and len(tok) == 4:
            return f"FY{tok[2:]}"
    return "FY?"


def financials_grid(c: Dict) -> List[List[str]]:
    """The 13x10 template table, filled from what the record actually holds
    (up to three CH-filed years). Everything else is 'n/d'. Pure, testable."""
    years = []  # oldest first: (label, rev, gp, pbt)
    for rk, dk in (("revenue_y3", "revenue_y3_date"), ("revenue_y2", "revenue_y2_date"),
                   ("revenue_y1", "revenue_y1_date")):
        if c.get(rk) is not None:
            years.append({"label": _yr(c.get(dk)), "rev": c.get(rk)})
    if years:
        years[-1]["gp"] = c.get("gross_profit_y1")
        years[-1]["pbt"] = c.get("profit_y1")

    def cell(i, key):
        return _gbp(years[i][key]) if i < len(years) and years[i].get(key) is not None else "n/d"

    def growth(i):
        if 0 < i < len(years):
            try:
                prev, cur = float(years[i - 1]["rev"]), float(years[i]["rev"])
                if prev:
                    return f"{(cur / prev - 1) * 100:+.1f}%"
            except (TypeError, ValueError):
                pass
        return "n/d" if i < len(years) else "-"

    def margin(i, key):
        if i < len(years) and years[i].get(key) is not None and years[i].get("rev"):
            try:
                return f"{float(years[i][key]) / float(years[i]['rev']) * 100:.1f}%"
            except (TypeError, ValueError, ZeroDivisionError):
                pass
        return "n/d" if i < len(years) else "-"

    n = 9  # value columns
    def row(label, fn):
        return [label] + [fn(i) if i < len(years) else "-" for i in range(n)]

    grid = [
        None,   # row 0: Actual/Forecast band - left as the template has it
        ["£"] + [years[i]["label"] if i < len(years) else "-" for i in range(n)],
        row("Revenue", lambda i: cell(i, "rev")),
        row("% Growth", growth),
        row("Cost of Sale", lambda i: "n/d"),
        row("Gross Profit", lambda i: cell(i, "gp")),
        row("% Margin", lambda i: margin(i, "gp")),
        row("Overheads", lambda i: "n/d"),
        row("EBITDA", lambda i: "n/d"),
        row("% Margin", lambda i: "n/d" if i < len(years) else "-"),
        row("Adjusted EBITDA", lambda i: "n/d"),
        row("PBT", lambda i: cell(i, "pbt")),
        ["Note: filled from Companies House filings held on the record; 'n/d' = not disclosed. "
         "No forecasts exist pre-diligence - the dataroom model fills this table later."]
        + [""] * n,
    ]
    return grid


# ── One grounded call for every written section ──────────────────────────────

_SCHEMA_HINT = """{
 "screening": {"tag": "one line, e.g. 'Proprietary outreach - UK vertical SaaS'", "rating": "Strong|Mixed|Weak", "note": "one sentence for the summary slide"},
 "business_description": [{"lead": "bold lead-in", "text": "rest of the bullet"}],
 "alignment": {"business_quality": {"verdict": "8-14 words", "tone": "good|warn|bad"},
               "revenue_quality": {...}, "market_dynamics": {...}, "management": {...},
               "deal_structuring": {...}, "sector_focus": {...}},
 "market_overview": [{"lead": "bold lead-in", "text": "rest; end researched lines with (AI research)"}],
 "thesis": ["3-5 bullet strings"],
 "value_creation": [{"lead": "...", "text": "..."}],
 "deal_structuring": ["3-5 bullet strings"],
 "risks": [{"category": "MARKET DYNAMICS|BUSINESS QUALITY|REVENUE QUALITY|MANAGEMENT|DEAL STRUCTURING", "risk": "...", "mitigant": "..."}],
 "exit": {"buyer_types": ["Strategic - ...", "Private equity - ..."], "examples": "one line of plausible named buyer examples (AI research)"},
 "questions": [{"lead": "topic, 2-4 words", "text": "the actual question"}]
}"""


def compose_deck(company: Dict, emails: List[dict], doc_summaries: List[str],
                 facts_block: str, email_block: str) -> Dict:
    """The written sections. Numbers only from the record; research is tagged."""
    api_key = os.getenv("GEMINI_API_KEY")
    empty = {"screening": {"tag": "", "rating": "Mixed", "note": ""},
             "business_description": [], "alignment": {}, "market_overview": [],
             "thesis": [], "value_creation": [], "deal_structuring": [],
             "risks": [], "exit": {"buyer_types": [], "examples": ""}, "questions": []}
    if not api_key:
        return empty
    docs = "\n".join(f"- {d}" for d in doc_summaries if d) or "(none on file)"
    prompt = f"""You are an associate at Averroes Capital preparing a 4-slide IC screening deck on "{company.get('name')}", an ORIGINATION-stage company (pre-diligence, sourced by our own founder outreach; there is no CIM and no dataroom).

HARD RULES:
- Company facts come ONLY from the record below. NEVER invent a number. Where a section normally needs undisclosed data, say "Not disclosed pre-diligence".
- You may use Google Search for MARKET and COMPETITOR context and plausible exit buyers. EVERY claim that comes from search rather than the record must end with "(AI research)".
- No em dashes anywhere. Be concise: bullets of one to two lines.
- 6-9 business_description bullets, 5-7 market_overview bullets, 6-10 risks, 8-12 questions specific to THIS company.
- Rating: Strong / Mixed / Weak against the mandate below.

THE MANDATE: UK/Ireland B2B software, roughly £2.5m-£10m revenue sweet spot, profitable or near, founder-led, buyout with founder liquidity.

COMPANY RECORD (the only permitted source for company facts):
{facts_block}

DOCUMENTS THE FOUNDER SENT US (AI summaries):
{docs}

EMAIL THREAD WITH THE FOUNDER (oldest first):
{email_block}

Return ONLY valid JSON exactly in this shape:
{_SCHEMA_HINT}"""
    try:
        from google import genai
        from google.genai.types import GenerateContentConfig, GoogleSearch, Tool
        client = genai.Client(api_key=api_key)
        response = client.models.generate_content(
            model="gemini-2.5-flash", contents=prompt,
            config=GenerateContentConfig(tools=[Tool(google_search=GoogleSearch())], temperature=0.2))
        text = (response.text or "").strip()
        if text.startswith("```"):
            text = text.split("\n", 1)[1].rsplit("```", 1)[0].strip()
        got = json.loads(text[text.find("{"):text.rfind("}") + 1])
        for k in empty:
            got.setdefault(k, empty[k])

        def _clean(v):
            if isinstance(v, str):
                return v.replace("—", "-").replace("–", "-")
            if isinstance(v, list):
                return [_clean(x) for x in v]
            if isinstance(v, dict):
                return {kk: _clean(vv) for kk, vv in v.items()}
            return v
        return _clean(got)
    except Exception as e:
        logger.warning(f"[ICDeck] compose failed for {company.get('name')}: {e}")
        return empty


# ── Filling the template ─────────────────────────────────────────────────────

def _shape(slide, shape_id):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape {shape_id} missing from template slide")


def _set_para_text(p, lead: str, text: str, proto_r):
    """Rebuild a paragraph's runs from a prototype run: bold lead + plain rest."""
    for r in list(p.findall(".//" + proto_r.tag)):
        p.remove(r)
    pieces = []
    if lead:
        pieces.append((lead + (" - " if text else ""), True))
    if text:
        pieces.append((text, False))
    if not pieces:
        pieces = [("", False)]
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for content, bold in pieces:
        r = copy.deepcopy(proto_r)
        rPr = r.find(ns + "rPr")
        if rPr is None:
            rPr = r.makeelement(ns + "rPr", {}); r.insert(0, rPr)
        rPr.set("b", "1" if bold else "0")
        t = r.find(ns + "t")
        if t is None:
            t = r.makeelement(ns + "t", {}); r.append(t)
        t.text = content
        p.append(r)


def _fill_bullets(shape, items):
    """Replace a shape's paragraphs with `items`, cloning the first paragraph's
    styling so the template's bullet marks, size and spacing survive."""
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    txBody = shape.text_frame._txBody
    paras = txBody.findall(ns + "p")
    proto = paras[0]
    proto_r = proto.find(ns + "r")
    if proto_r is None:  # empty prototype: fabricate a minimal run
        proto_r = proto.makeelement(ns + "r", {})
        t = proto_r.makeelement(ns + "t", {}); t.text = ""
        proto_r.append(t)
    for p in paras:
        txBody.remove(p)
    for it in items:
        lead, text = (it.get("lead", ""), it.get("text", "")) if isinstance(it, dict) else ("", str(it))
        p = copy.deepcopy(proto)
        for r in list(p.findall(ns + "r")):
            p.remove(r)
        _set_para_text(p, lead, text, proto_r)
        txBody.append(p)


def _set_text(shape, text: str):
    """Replace a shape's text, keeping the first run's formatting."""
    _fill_bullets(shape, [{"lead": "", "text": text}])


def _set_cell(cell, text: str):
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    txBody = cell.text_frame._txBody
    paras = txBody.findall(ns + "p")
    proto = paras[0]
    proto_r = proto.find(ns + "r")
    if proto_r is None:
        proto_r = proto.makeelement(ns + "r", {})
        t = proto_r.makeelement(ns + "t", {}); t.text = ""
        proto_r.append(t)
    for p in paras[1:]:
        txBody.remove(p)
    for r in list(proto.findall(ns + "r")):
        proto.remove(r)
    _set_para_text(proto, "", text, proto_r)


def render_deck(company: Dict, content: Dict, grid: List[List[str]]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    prs = Presentation(TEMPLATE_PATH)
    name = company.get("name") or "Company"
    s1, s2, s3, s4 = prs.slides

    # ── Slide 1: screening summary ──
    scr = content.get("screening") or {}
    rating = scr.get("rating") if scr.get("rating") in _RATING_FILL else "Mixed"
    rev = company.get("revenue_y1")
    rev_line = (f"{_gbp(rev)} (CH filed)" if rev is not None
                else (f"£{company.get('revenue_estimate_m'):.1f}m (est.)"
                      if company.get("revenue_estimate_m") is not None else "n/d"))
    _set_text(_shape(s1, 11), name)
    _set_text(_shape(s1, 12), scr.get("tag") or f"Proprietary outreach - {company.get('sector') or 'UK B2B software'}")
    _set_text(_shape(s1, 13), rev_line)
    _set_text(_shape(s1, 15), rating)
    chip = _shape(s1, 14)
    chip.fill.solid(); chip.fill.fore_color.rgb = RGBColor(*_RATING_FILL[rating])
    _set_text(_shape(s1, 16), scr.get("note") or "Origination-stage screen: sourced by AverroesIntel founder outreach; no CIM, pre-diligence.")

    # ── Slide 2: overview ──
    _set_text(_shape(s2, 3), name)
    _fill_bullets(_shape(s2, 7), content.get("business_description") or [{"lead": "", "text": "Not disclosed pre-diligence."}])
    _fill_bullets(_shape(s2, 9), content.get("market_overview") or [{"lead": "", "text": "Not disclosed pre-diligence."}])
    chip_text_ids = [27, 29, 31, 33, 35, 37]
    chip_bg_ids = [26, 28, 30, 32, 34, 36]
    align = content.get("alignment") or {}
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for (key, label), tid, bid in zip(ALIGNMENT_KEYS, chip_text_ids, chip_bg_ids):
        got = align.get(key) or {}
        sh = _shape(s2, tid)
        txBody = sh.text_frame._txBody
        paras = txBody.findall(ns + "p")
        # paragraph 0 = the fixed label, paragraph 1+ = the verdict
        proto = paras[1] if len(paras) > 1 else paras[0]
        proto_r = proto.find(ns + "r")
        for p in paras[1:]:
            txBody.remove(p)
        if proto_r is None:
            proto_r = proto.makeelement(ns + "r", {})
            t = proto_r.makeelement(ns + "t", {}); t.text = ""
            proto_r.append(t)
        p = copy.deepcopy(proto)
        for r in list(p.findall(ns + "r")):
            p.remove(r)
        _set_para_text(p, "", got.get("verdict") or "Not assessed", proto_r)
        txBody.append(p)
        bg = _shape(s2, bid)
        tone = got.get("tone") if got.get("tone") in _TONE_FILL else "warn"
        bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(*_TONE_FILL[tone])
    tbl = _shape(s2, 40).table
    for ri, row in enumerate(grid):
        if row is None:
            continue
        for ci in range(min(len(row), len(tbl.columns))):
            _set_cell(tbl.cell(ri, ci), row[ci])

    # ── Slide 3: rationale ──
    _set_text(_shape(s3, 3), name)
    _fill_bullets(_shape(s3, 7), content.get("thesis") or ["Not yet formed - pre-diligence."])
    _fill_bullets(_shape(s3, 9), content.get("value_creation") or [{"lead": "", "text": "Hypotheses to be formed after the first call."}])
    _fill_bullets(_shape(s3, 24), content.get("deal_structuring") or ["Not yet discussed with the founder."])
    risks = (content.get("risks") or [])[:10]
    rt = _shape(s3, 26).table
    for ri in range(1, len(rt.rows)):
        if ri - 1 < len(risks):
            r = risks[ri - 1]
            _set_cell(rt.cell(ri, 0), (r.get("category") or "").upper())
            _set_cell(rt.cell(ri, 1), r.get("risk") or "")
            _set_cell(rt.cell(ri, 2), r.get("mitigant") or "")
        else:
            for ci in range(3):
                _set_cell(rt.cell(ri, ci), "")
    exit_ = content.get("exit") or {}
    exit_items = [{"lead": "Buyer types", "text": "; ".join(exit_.get("buyer_types") or []) or "Not yet assessed"}]
    if exit_.get("examples"):
        exit_items.append({"lead": "Examples", "text": exit_["examples"]})
    _fill_bullets(_shape(s3, 11), exit_items)

    # ── Slide 4: questions ──
    _set_text(_shape(s4, 3), name)
    qs = content.get("questions") or []
    # The template shape numbers its paragraphs itself (buAutoNum), so the
    # lead must NOT carry a manual "1." - it would render as "1. 1. Topic".
    _fill_bullets(_shape(s4, 6),
                  [{"lead": q.get("lead", ""), "text": q.get("text", "")} for q in qs]
                  or ["No questions generated."])

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
